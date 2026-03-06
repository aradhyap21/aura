"""
AURA v3 — Document Extractor
Handles:
  - PDF (normal text-based)
  - PDF (scanned / handwritten) → OCR via Tesseract
  - DOCX
  - PPTX

Returns clean full text for the semantic topic extractor.
"""

import io
import re
import pdfplumber
from docx import Document
from pptx import Presentation


# ─────────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────────

def clean_text(text: str) -> str:
    """Remove noise common in extracted documents."""
    text = re.sub(r'^\s*\d{1,3}\s*$', '', text, flags=re.MULTILINE)
    text = re.sub(r'https?://\S+', '', text)
    text = re.sub(r'\n{4,}', '\n\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'[^\x20-\x7E\n]', ' ', text)
    return text.strip()


def is_scanned_pdf(file_bytes: bytes) -> bool:
    """
    Detect if a PDF is scanned (image-only, no text layer).
    If less than 50 words extracted from first 3 pages = scanned.
    """
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            sample = ""
            for page in pdf.pages[:3]:
                t = page.extract_text()
                if t:
                    sample += t
            return len(sample.split()) < 50
    except Exception:
        return False


# ─────────────────────────────────────────────────────────────────
# PDF
# ─────────────────────────────────────────────────────────────────

def extract_pdf_text(file_bytes: bytes) -> str:
    pages = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                pages.append(t)
    return clean_text('\n\n'.join(pages))


def extract_pdf_ocr(file_bytes: bytes) -> str:
    """
    OCR for scanned/handwritten PDFs.
    Requires: Tesseract + pytesseract + pdf2image + Pillow

    Install Tesseract:
      Windows: https://github.com/UB-Mannheim/tesseract/wiki
      Linux:   sudo apt install tesseract-ocr
    """
    try:
        import pytesseract
        from pdf2image import convert_from_bytes

        print("  Scanned PDF detected — running OCR...")
        images = convert_from_bytes(file_bytes, dpi=300)
        pages  = []
        for i, img in enumerate(images):
            print(f"  OCR page {i+1}/{len(images)}...", end="\r")
            text = pytesseract.image_to_string(img, lang='eng')
            if text.strip():
                pages.append(text)
        print()
        return clean_text('\n\n'.join(pages))

    except ImportError as e:
        raise ValueError(
            f"Scanned PDF detected but OCR not installed.\n"
            f"pip install pytesseract pdf2image Pillow\n"
            f"Tesseract: https://github.com/UB-Mannheim/tesseract/wiki\n"
            f"Error: {e}"
        )


def extract_pdf(file_bytes: bytes) -> str:
    if is_scanned_pdf(file_bytes):
        return extract_pdf_ocr(file_bytes)
    text = extract_pdf_text(file_bytes)
    if len(text.split()) < 30:
        raise ValueError("Very little text extracted. May be a scanned PDF.")
    return text


# ─────────────────────────────────────────────────────────────────
# DOCX
# ─────────────────────────────────────────────────────────────────

def extract_docx(file_bytes: bytes) -> str:
    doc   = Document(io.BytesIO(file_bytes))
    lines = []

    for para in doc.paragraphs:
        text  = para.text.strip()
        style = para.style.name.lower() if para.style else ""
        if not text:
            continue
        if "heading 1" in style:
            lines.append(f"# {text}")
        elif "heading 2" in style:
            lines.append(f"## {text}")
        elif "heading 3" in style or "heading 4" in style:
            lines.append(f"### {text}")
        else:
            lines.append(text)

    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                lines.append(row_text)

    text = clean_text('\n'.join(lines))
    if len(text.split()) < 30:
        raise ValueError("DOCX appears empty.")
    return text


# ─────────────────────────────────────────────────────────────────
# PPTX
# ─────────────────────────────────────────────────────────────────

def extract_pptx(file_bytes: bytes) -> str:
    prs   = Presentation(io.BytesIO(file_bytes))
    lines = []

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n--- Slide {i} ---")

        if slide.shapes.title and slide.shapes.title.text.strip():
            lines.append(f"# {slide.shapes.title.text.strip()}")

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text and len(text) > 3:
                    lines.append(text)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes and len(notes.split()) > 5:
                lines.append(f"[Notes: {notes}]")

    text = clean_text('\n'.join(lines))
    if len(text.split()) < 20:
        raise ValueError("PPTX appears empty.")
    return text


# ─────────────────────────────────────────────────────────────────
# ROUTER
# ─────────────────────────────────────────────────────────────────

def detect_file_type(content_type: str, filename: str) -> str:
    ct = (content_type or "").lower()
    fn = (filename or "").lower()
    if "pdf" in ct or fn.endswith(".pdf"):
        return "pdf"
    if "wordprocessingml" in ct or fn.endswith(".docx"):
        return "docx"
    if "presentationml" in ct or fn.endswith((".pptx", ".ppt")):
        return "pptx"
    raise ValueError(f"Unsupported file type: {content_type or filename}")


def extract_full_text(file_bytes: bytes, file_type: str) -> str:
    if file_type == "pdf":
        return extract_pdf(file_bytes)
    elif file_type == "docx":
        return extract_docx(file_bytes)
    elif file_type == "pptx":
        return extract_pptx(file_bytes)
    raise ValueError(f"Unknown file type: {file_type}")